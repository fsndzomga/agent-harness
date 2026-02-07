"""Custom tools for GAIA Agent v1.

Includes: smart search, Wikipedia, memory (save/read notes),
enhanced file reading (PDF, Excel, CSV, audio, images, docx, json, html).
"""

from __future__ import annotations

import json
import os
import re
import sys
import tempfile
from typing import Any

from smolagents import tool

# ---------------------------------------------------------------------------
# Module-level memory state (set per-task by runner.py)
# ---------------------------------------------------------------------------
_memory_dir: str | None = None


def set_memory_dir(path: str) -> None:
    """Initialise the per-task memory directory."""
    global _memory_dir
    _memory_dir = path
    os.makedirs(path, exist_ok=True)
    notes_path = os.path.join(path, "notes.json")
    if not os.path.exists(notes_path):
        with open(notes_path, "w") as f:
            json.dump({}, f)


def get_memory_dir() -> str:
    assert _memory_dir is not None, "Memory not initialised — call set_memory_dir first"
    return _memory_dir


# ===================================================================
# 1.  Smart web search (DuckDuckGo + reformulation + retry)
# ===================================================================

def _reformulate(query: str) -> list[str]:
    """Generate 1–2 alternative search queries."""
    alts: list[str] = []
    q = query.strip()
    words = q.split()

    # Simplified version: just try with fewer words
    if len(words) > 6:
        alts.append(" ".join(words[:6]))

    return alts[:2]


@tool
def search_web(query: str) -> str:
    """Search the web using DuckDuckGo. If few results are found, automatically
    retries with a simplified query.

    Args:
        query: The search query — be specific and include key terms.

    Returns:
        Formatted search results with titles, URLs, and snippets.
    """
    from duckduckgo_search import DDGS

    all_results: list[dict] = []

    # --- first attempt ---
    try:
        with DDGS() as ddgs:
            all_results.extend(ddgs.text(query, max_results=10))
    except Exception:
        pass

    # --- retry with simplification if no results ---
    if len(all_results) < 2:
        for alt in _reformulate(query):
            try:
                with DDGS() as ddgs:
                    all_results.extend(ddgs.text(alt, max_results=8))
            except Exception:
                pass
            if len(all_results) >= 3:
                break

    # --- deduplicate by URL ---
    seen: set[str] = set()
    unique: list[dict] = []
    for r in all_results:
        url = r.get("href", "")
        if url and url not in seen:
            seen.add(url)
            unique.append(r)

    if not unique:
        return f"No search results found for: {query}. Try a different query."

    lines = [f"Search results for: {query}", ""]

    for i, r in enumerate(unique[:10], 1):
        title = r.get("title", "N/A")
        url = r.get("href", "N/A")
        body = r.get("body", "")[:300]
        lines.append(f"{i}. {title}")
        lines.append(f"   URL: {url}")
        lines.append(f"   {body}")
        lines.append("")

    return "\n".join(lines)


# ===================================================================
# 2.  Wikipedia search
# ===================================================================

@tool
def wiki_search(query: str) -> str:
    """Search Wikipedia for factual / encyclopaedic information.

    Args:
        query: The topic to look up on Wikipedia.

    Returns:
        Article summary plus links to related articles.
    """
    import urllib.request
    import urllib.parse

    params = urllib.parse.urlencode({
        "action": "query",
        "list": "search",
        "srsearch": query,
        "format": "json",
        "srlimit": 3,
    })
    search_url = f"https://en.wikipedia.org/w/api.php?{params}"

    try:
        with urllib.request.urlopen(search_url, timeout=10) as resp:
            data = json.loads(resp.read().decode())

        hits = data.get("query", {}).get("search", [])
        if not hits:
            return f"No Wikipedia results for: {query}"

        title = hits[0]["title"]
        summary_url = (
            "https://en.wikipedia.org/api/rest_v1/page/summary/"
            + urllib.parse.quote(title)
        )

        with urllib.request.urlopen(summary_url, timeout=10) as resp:
            sdata = json.loads(resp.read().decode())

        extract = sdata.get("extract", "No summary available.")
        page_url = (
            sdata.get("content_urls", {}).get("desktop", {}).get("page", "")
        )

        out = f"# {title}\n\n{extract}\n\nSource: {page_url}"

        if len(hits) > 1:
            out += "\n\n## Related articles:"
            for h in hits[1:]:
                snippet = re.sub(r"<[^>]+>", "", h.get("snippet", ""))[:120]
                out += f"\n- {h['title']}: {snippet}"

        return out
    except Exception as e:
        return f"Wikipedia search error: {e}"


# ===================================================================
# 3.  Persistent memory (notes)
# ===================================================================

@tool
def save_note(key: str, content: str) -> str:
    """Save a note to persistent memory for later retrieval.

    Args:
        key: Short descriptive key (e.g. "population", "table_data").
        content: The content to save.

    Returns:
        Confirmation message.
    """
    notes_path = os.path.join(get_memory_dir(), "notes.json")
    with open(notes_path, "r") as f:
        notes = json.load(f)
    notes[key] = content
    with open(notes_path, "w") as f:
        json.dump(notes, f, indent=2)
    return f"Saved note '{key}' ({len(content)} chars). Total notes: {len(notes)}"


@tool
def read_notes() -> str:
    """Read all saved notes from persistent memory.

    Returns:
        All saved notes.
    """
    notes_path = os.path.join(get_memory_dir(), "notes.json")
    with open(notes_path, "r") as f:
        notes = json.load(f)
    if not notes:
        return "No notes saved yet."
    parts = []
    for k, v in notes.items():
        parts.append(f"## {k}\n{v}")
    return "\n\n".join(parts)


# ===================================================================
# 4.  Enhanced file reading
# ===================================================================

@tool
def read_file_tool(file_path: str) -> str:
    """Read and extract text from a file. Supports PDF, Excel, CSV, JSON,
    HTML, DOCX, PPTX, images (OCR), audio (transcription), and plain text.

    Args:
        file_path: Path to the file to read.

    Returns:
        Extracted text content.
    """
    if not os.path.exists(file_path):
        return f"Error: File not found: {file_path}"

    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            return _read_pdf(file_path)
        elif ext in (".xlsx", ".xls"):
            return _read_excel(file_path)
        elif ext == ".csv":
            return _read_csv(file_path)
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"):
            return _read_image(file_path)
        elif ext in (".mp3", ".wav", ".m4a", ".ogg", ".flac"):
            return _read_audio(file_path)
        elif ext == ".json":
            return _read_json(file_path)
        elif ext in (".html", ".htm"):
            return _read_html(file_path)
        elif ext == ".xml":
            return _read_xml(file_path)
        elif ext == ".docx":
            return _read_docx(file_path)
        elif ext == ".pptx":
            return _read_pptx(file_path)
        else:
            with open(file_path, "r", errors="replace") as f:
                return f.read()[:100_000]
    except Exception as e:
        return f"Error reading {file_path}: {e}"


# ---- private readers -------------------------------------------------

def _read_pdf(path: str) -> str:
    import pymupdf
    doc = pymupdf.open(path)
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n".join(pages)[:100_000]


def _read_excel(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f"--- Sheet: {name} ---")
        for row in ws.iter_rows(values_only=True):
            parts.append("\t".join(str(c) if c is not None else "" for c in row))
    wb.close()
    return "\n".join(parts)[:100_000]


def _read_csv(path: str) -> str:
    with open(path, "r", errors="replace") as f:
        return f.read()[:100_000]


def _read_json(path: str) -> str:
    with open(path, "r") as f:
        data = json.load(f)
    return json.dumps(data, indent=2)[:100_000]


def _read_html(path: str) -> str:
    try:
        from markdownify import markdownify
        with open(path, "r", errors="replace") as f:
            html = f.read()
        return markdownify(html)[:100_000]
    except ImportError:
        with open(path, "r", errors="replace") as f:
            text = f.read()
        return re.sub(r"<[^>]+>", " ", text)[:100_000]


def _read_xml(path: str) -> str:
    with open(path, "r", errors="replace") as f:
        return f.read()[:100_000]


def _read_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.append("\t".join(c.text for c in row.cells))
        return "\n".join(parts)[:100_000]
    except ImportError:
        return "[DOCX file — install python-docx to read]"


def _read_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
        return "\n".join(parts)[:100_000]
    except ImportError:
        return "[PPTX file — install python-pptx to read]"


def _read_image(path: str) -> str:
    """Attempt OCR on image, fall back to description."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(path)
        text = pytesseract.image_to_string(img)
        if text.strip():
            return f"[OCR from image: {path}]\n{text}"
        return (
            f"[Image file: {path}, size={img.size}, mode={img.mode}] — "
            "OCR returned no text. Use code to analyse the image if needed."
        )
    except ImportError:
        try:
            from PIL import Image
            img = Image.open(path)
            return (
                f"[Image file: {path}, size={img.size}, mode={img.mode}] — "
                "pytesseract not installed for OCR. Use code to process."
            )
        except Exception:
            return f"[Image file: {path}] — Use code to process this image."


def _read_audio(path: str) -> str:
    """Attempt audio transcription via speech_recognition, fall back gracefully."""
    # --- method 1: speech_recognition + pydub ---
    try:
        import speech_recognition as sr  # type: ignore
        from pydub import AudioSegment    # type: ignore

        ext = os.path.splitext(path)[1].lower()
        wav_path = path
        if ext != ".wav":
            audio = AudioSegment.from_file(path)
            wav_path = os.path.join(tempfile.gettempdir(), "gaia_tmp_audio.wav")
            audio.export(wav_path, format="wav")

        recognizer = sr.Recognizer()
        with sr.AudioFile(wav_path) as source:
            audio_data = recognizer.record(source)

        text = recognizer.recognize_google(audio_data)
        return f"[Audio transcription of {path}]\n{text}"
    except ImportError:
        pass
    except Exception as e:
        print(f"[audio] speech_recognition failed: {e}", file=sys.stderr, flush=True)

    # --- method 2: whisper (local) ---
    try:
        import whisper  # type: ignore
        model = whisper.load_model("base")
        result = model.transcribe(path)
        return f"[Audio transcription of {path}]\n{result['text']}"
    except ImportError:
        pass
    except Exception as e:
        print(f"[audio] whisper failed: {e}", file=sys.stderr, flush=True)

    return (
        f"[Audio file: {path}] — Transcription libraries not available. "
        "Install SpeechRecognition + pydub, or openai-whisper."
    )


# ===================================================================
# 5.  Enhanced webpage visitor
# ===================================================================

@tool
def visit_webpage(url: str) -> str:
    """Fetch a webpage and return its content as clean Markdown text.

    Args:
        url: The URL to visit.

    Returns:
        The page content converted to Markdown, or an error message.
    """
    import urllib.request
    import urllib.error

    try:
        req = urllib.request.Request(
            url,
            headers={
                "User-Agent": (
                    "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 "
                    "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
                ),
                "Accept": "text/html,application/xhtml+xml,*/*",
                "Accept-Language": "en-US,en;q=0.9",
            },
        )
        with urllib.request.urlopen(req, timeout=20) as resp:
            content_type = resp.headers.get("Content-Type", "")
            raw = resp.read()

        # Decode
        charset = "utf-8"
        if "charset=" in content_type:
            charset = content_type.split("charset=")[-1].split(";")[0].strip()
        html = raw.decode(charset, errors="replace")

        # Convert to markdown
        try:
            from markdownify import markdownify
            text = markdownify(html)
        except ImportError:
            text = re.sub(r"<[^>]+>", " ", html)

        # Clean up excessive whitespace
        text = re.sub(r"\n{3,}", "\n\n", text)

        return text[:80_000]
    except urllib.error.HTTPError as e:
        return f"HTTP error {e.code}: {e.reason} — {url}"
    except urllib.error.URLError as e:
        return f"URL error: {e.reason} — {url}"
    except Exception as e:
        return f"Error visiting {url}: {e}"
